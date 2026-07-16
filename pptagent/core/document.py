"""PPTDocument in-memory data model.

This module defines the core data structures that represent a PowerPoint file
in memory: TextFormat, Element, Slide, PPTDocument, and related types.

All tools read from and write to instances of these classes.
"""

from __future__ import annotations

import io
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

EMU_PER_INCH: int = 914400
EMU_PER_CM: int = 360000


def _rgb_to_hex(rgb: Any) -> str:
    """Convert a python-pptx RGBColor to a '#RRGGBB' hex string."""
    try:
        return f"#{rgb!s}"
    except Exception:
        return "#000000"


def _classify_shape(shape: Any) -> str:
    """Return the semantic element type for python-pptx *shape*."""

    from pptx.enum.shapes import MSO_SHAPE_TYPE

    stype = shape.shape_type

    # Does this shape carry visible text?
    has_text = bool(
        hasattr(shape, "has_text_frame")
        and shape.has_text_frame
        and shape.text_frame.text.strip()
    )

    if stype == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if stype == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if stype == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if stype == MSO_SHAPE_TYPE.DIAGRAM:
        return "diagram"
    if stype == MSO_SHAPE_TYPE.MEDIA:
        return "media"
    if stype == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if stype == MSO_SHAPE_TYPE.LINKED_OLE_OBJECT:
        # Check ProgID for formula-like OLE objects.
        try:
            if hasattr(shape, "ole_format") and shape.ole_format is not None:
                prog_id = getattr(shape.ole_format, "prog_id", "") or ""
                if prog_id.lower() in {
                    "equation",
                    "equations",
                    "equations.3",
                    "mathtype",
                    "mathtype.equation",
                    "formula",
                    "chemdraw",
                    "chemdraw.document",
                    "math",
                }:
                    return "formula"
        except Exception:
            pass
        return "shape"
    if stype == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "text" if has_text else "shape"
    if stype == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "text" if has_text else "shape"
    if has_text:
        return "text"
    return "shape"


# Extract content from a single shape.
def _extract_text(shape: Any) -> str:
    """Return all text inside *shape*'s text frame, or ''."""
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        return str(shape.text_frame.text)
    return ""


def _extract_image_content(shape: Any) -> dict[str, Any]:
    """Build the ``content`` dict for an image shape."""
    try:
        image = shape.image
        blob: bytes = image.blob
        content_type: str = image.content_type or "image/png"
        ext: str = getattr(image, "ext", None) or content_type.split("/")[-1] or "png"
        filename = f"image_{uuid.uuid4().hex[:8]}.{ext}"
        return {"blob": blob, "content_type": content_type, "filename": filename}
    except Exception:
        return {"blob": b"", "content_type": "image/png", "filename": "unknown.png"}


def _extract_table_content(shape: Any) -> list[list[str]]:
    """Return table cell text as a 2-D list of strings."""
    content: list[list[str]] = []
    try:
        for row in shape.table.rows:
            content.append([cell.text for cell in row.cells])
    except Exception:
        pass
    return content


def _extract_chart_content(shape: Any) -> dict[str, Any]:
    """Return chart metadata and series data."""
    try:
        chart = shape.chart
        chart_type = str(chart.chart_type) if chart.chart_type is not None else "unknown"
        data: dict[str, Any] = {"chart_type": chart_type}

        # Series
        series_list: list[dict[str, Any]] = []
        for ser in chart.series:
            info: dict[str, Any] = {"values": list(ser.values)}
            if hasattr(ser, "name"):
                info["name"] = ser.name
            series_list.append(info)
        data["series"] = series_list

        # Categories from the first plot
        if chart.plots:
            plot = chart.plots[0]
            if hasattr(plot, "categories"):
                data["categories"] = list(plot.categories)
        return data
    except Exception:
        return {"chart_type": "unknown", "error": "Could not extract chart data"}


def _extract_group_shapes(
    group_shape: Any, z_start: int
) -> list[tuple[Any, int]]:
    """Flatten one level of group nesting.

    Returns ``[(shape, z_order), ...]`` so the caller can continue processing.
    """
    result: list[tuple[Any, int]] = []
    z = z_start
    try:
        for sub in group_shape.shapes:
            if _classify_shape(sub) == "group":
                nested = _extract_group_shapes(sub, z)
                result.extend(nested)
                z += len(nested)
            else:
                result.append((sub, z))
                z += 1
    except Exception:
        pass
    return result


def _extract_text_format(shape: Any) -> TextFormat | None:
    """Build a TextFormat from the first run of the first paragraph."""
    try:
        if not (hasattr(shape, "has_text_frame") and shape.has_text_frame):
            return None
        tf = shape.text_frame
        if not tf.paragraphs:
            return None
        para = tf.paragraphs[0]
        if not para.runs:
            return None

        fmt = TextFormat.from_run(para.runs[0])

        # Paragraph-level alignment
        from pptx.enum.text import PP_ALIGN

        _align_map = {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "justify",
        }
        if para.alignment is not None:
            fmt.alignment = _align_map.get(para.alignment, "left")

        return fmt
    except Exception:
        return None


def _build_position(shape: Any) -> tuple[float, float, float, float]:
    """Return ``(left, top, width, height)`` in EMU, defaulting missing values to 0."""
    return (
        float(getattr(shape, "left", 0) or 0),
        float(getattr(shape, "top", 0) or 0),
        float(getattr(shape, "width", 0) or 0),
        float(getattr(shape, "height", 0) or 0),
    )


@dataclass
class TextFormat:
    """Complete description of text formatting properties.

    Its goals:
    - Mirrors full-set properties in OOXML CT_TextCharacterProperties.
    - All colour values are stored as HEX strings.
    """

    font_name: str = "Calibri"
    font_size: float = 12.0  # points (pt)
    font_color: str = "#000000"  # HEX RGBA or "theme"
    bold: bool = False
    italic: bool = False
    underline: bool = False
    strikethrough: bool = False
    superscript: bool = False
    subscript: bool = False
    highlight: str | None = None  # highlight colour HEX, or None
    alignment: str = "left"  # "left" | "center" | "right" | "justify"
    line_spacing: float = 1.0  # multiplier (1.0 = single)
    paragraph_spacing: float = 0.0  # space before/after in pt

    @classmethod
    def from_run(cls, run: Any) -> "TextFormat":
        """Extract TextFormat from a python-pptx Run object.

        Reads font properties from the run's ``font`` attribute.  All colour
        values are normalised to ``#RRGGBB`` hex strings.

        Returns a ``TextFormat`` with defaults for any properties that cannot
        be read.
        """
        try:
            font = run.font
        except (AttributeError, TypeError):
            return cls()

        font_name: str = font.name or "Calibri"
        font_size: float = 12.0
        raw_size = font.size
        if raw_size is not None:
            try:
                sz = float(raw_size)
                # python-pptx may return centipoints (>100), EMU (>100000), or points
                if sz > 100_000:  # EMU  (12 pt = 152 400 EMU)
                    font_size = sz / 12_700.0
                elif sz > 100:  # centipoints  (12 pt = 1_200)
                    font_size = sz / 100.0
                else:  # already in points
                    font_size = sz
            except (TypeError, ValueError):
                pass

        font_color: str = "#000000"
        try:
            if font.color and font.color.rgb is not None:
                font_color = _rgb_to_hex(font.color.rgb)
            elif font.color and font.color.theme_color is not None:
                font_color = "theme"
        except Exception:
            pass

        bold: bool = bool(font.bold) if font.bold is not None else False
        italic: bool = bool(font.italic) if font.italic is not None else False
        underline: bool = bool(font.underline) if font.underline is not None else False

        strikethrough: bool = False
        for attr in ("strikethrough", "strike"):
            val = getattr(font, attr, None)
            if val is not None:
                strikethrough = bool(val)
                break

        superscript: bool = False
        subscript: bool = False

        return cls(
            font_name=font_name,
            font_size=font_size,
            font_color=font_color,
            bold=bold,
            italic=italic,
            underline=underline,
            strikethrough=strikethrough,
            superscript=superscript,
            subscript=subscript,
        )


@dataclass
class Element:
    """A single foreground element on a slide.

    This is an abstraction over python-pptx shapes,
    grouping them by semantic type.
    """

    id: str  # unique id, format: "{type}_{uuid8}"
    type: str  # "text" | "image" | "table" | "chart" | "formula" | "diagram" | "media" | "shape"
    position: tuple[float, float, float, float]  # (left, top, width, height) in EMU
    content: Any  # type-specific payload (str, bytes, etc.)
    format: TextFormat | None = None  # text formatting; None for non-text elements
    z_order: int = 0  # stacking order
    metadata: dict[str, Any] = field(default_factory=dict)  # animation, hyperlinks, alt-text, ...

    # static EMU conversion helpers
    @staticmethod
    def emu_to_inches(emu: int) -> float:
        return emu / EMU_PER_INCH

    @staticmethod
    def inches_to_emu(inches: float) -> int:
        return int(inches * EMU_PER_INCH)

    @staticmethod
    def emu_to_cm(emu: int) -> float:
        return emu / EMU_PER_CM

    @staticmethod
    def cm_to_emu(cm: float) -> int:
        return int(cm * EMU_PER_CM)


@dataclass
class Slide:
    """A single slide within a ppt."""

    index: int  # 0-based position in the slide list
    layout_name: str = "Blank"  # e.g. "Title Slide"
    elements: dict[str, Element] = field(default_factory=dict)  # element_id -> Element
    background: dict[str, Any] = field(default_factory=dict)  # fill colour, image, ...
    notes: str = ""  # speaker notes (plain text)


@dataclass
class CursorState:
    """Tracks the user's current position in the document."""

    current_slide_index: int = 0
    selected_element_ids: list[str] = field(default_factory=list)


@dataclass
class BackgroundInfo:
    """Background fill for a slide."""

    fill_color: str | None = None  # HEX
    image_path: str | None = None


@dataclass
class PPTDocument:
    """In-memory representation of a PowerPoint file.

    **This is the central data structure that all tools read from and write to**.
    It provides an element-oriented view of content,
    wrapping a ``python-pptx`` ``Presentation`` object.
    """

    file_path: str | None = None  # None when created from scratch (not yet saved)
    file_name: str = "Untitled.pptx"
    slides: list[Slide] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)

    # ChangeManager inits lazily via record_change()
    _change_manager: Any = field(default=None, repr=False)

    cursor: CursorState = field(default_factory=CursorState)

    # Factory methods
    @classmethod
    def from_file(cls, path: str | Path) -> "PPTDocument":
        """Load a PPTDocument from a ``.pptx`` file using python-pptx.

        Iterates over every slide and every shape,
        classifying each shape into an `Element` type,
        and extracting its content, position, and formatting.
        """
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required to open .pptx files.  "
                "Install it with: pip install python-pptx"
            ) from exc

        path = Path(path).resolve()
        if not path.exists():
            raise FileNotFoundError(f"PPT file not found: {path}")

        prs = Presentation(str(path))

        doc = cls(
            file_path=str(path),
            file_name=path.name,
            slides=[],
        )

        for slide_idx, pptx_slide in enumerate(prs.slides):
            slide = Slide(
                index=slide_idx,
                layout_name=pptx_slide.slide_layout.name
                if pptx_slide.slide_layout
                else "Blank",
            )

            # Collect shapes, flattening one level of groups
            shapes_to_process: list[tuple[Any, int]] = []
            z = 0
            for shape in pptx_slide.shapes:
                if _classify_shape(shape) == "group":
                    nested = _extract_group_shapes(shape, z)
                    shapes_to_process.extend(nested)
                    z += len(nested)
                else:
                    shapes_to_process.append((shape, z))
                    z += 1

            # Build elements from shapes
            for shape, z_order in shapes_to_process:
                element = _shape_to_element(shape, z_order)
                if element is not None:
                    slide.elements[element.id] = element

            doc.slides.append(slide)

        return doc

    @classmethod
    def from_scratch(cls) -> "PPTDocument":
        """Create an empty single-slide ppt."""
        slide = Slide(index=0, layout_name="Blank")
        return cls(slides=[slide])

    # Serialize
    def to_file(self, path: str | Path | None = None) -> str:
        """Write the document back to a ``.pptx`` file.

        If *path* is ``None``, the document is written to ``self.file_path``.

        Basic element-roundtrip support for **text** and **image** elements;
        other element types are silently skipped (Phase 2 scope).

        Returns the absolute path of the written file.
        """
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required to write .pptx files.  "
                "Install it with: pip install python-pptx"
            ) from exc

        save_path = Path(path) if path else Path(self.file_path or "")
        if not save_path.name:
            raise ValueError(
                "No output path specified and PPTDocument.file_path is None"
            )

        prs = Presentation()

        for slide_data in self.slides:
            # Pick a blank layout or fall back to the first available one.
            blank_layout = None
            for layout in prs.slide_layouts:
                if "blank" in layout.name.lower():
                    blank_layout = layout
                    break
            if blank_layout is None:
                blank_layout = prs.slide_layouts[0]

            pptx_slide = prs.slides.add_slide(blank_layout)

            for element in slide_data.elements.values():
                left, top, width, height = element.position

                if element.type == "text":
                    _write_text_element(pptx_slide, element, left, top, width, height)

                elif element.type == "image":
                    _write_image_element(pptx_slide, element, left, top, width, height)

                elif element.type == "table":
                    _write_table_element(pptx_slide, element, left, top, width, height)

                elif element.type == "chart":
                    _write_chart_element(pptx_slide, element, left, top, width, height)

        save_path = save_path.resolve()
        prs.save(str(save_path))
        self.file_path = str(save_path)
        self.file_name = save_path.name
        return str(save_path)

    # Slide access
    def get_slide(self, index: int) -> Slide:
        """Return the slide at *index*, 0-based."""
        if not 0 <= index < len(self.slides):
            raise IndexError(
                f"Slide index {index} out of range (0 -{len(self.slides) - 1})"
            )
        return self.slides[index]

    def get_element(self, slide_index: int, element_id: str) -> Element:
        """Return the element with *element_id* on the given slide."""
        slide = self.get_slide(slide_index)
        if element_id not in slide.elements:
            raise KeyError(
                f"Element '{element_id}' not found on slide {slide_index}"
            )
        return slide.elements[element_id]

    # Undo / redo delegation
    def record_change(self, change: Any) -> None:
        """Record an undoable change.

        If a `ChangeManager` has not been assigned yet,
        this method raises `RuntimeError`.
        To create a `ChangeManager` for a document,
        use `self.set_change_manager` first.
        """
        if self._change_manager is None:
            raise RuntimeError(
                "No ChangeManager attached to this document. "
                "Call set_change_manager() first."
            )
        self._change_manager.record(change)

    def set_change_manager(self, cm: Any) -> None:
        """Attach a ChangeManager instance to this document."""
        self._change_manager = cm

    @property
    def change_manager(self) -> Any:
        """Return the attached ChangeManager, may be None."""
        return self._change_manager

    def undo(self) -> Any:
        """Undo the most recent change. Delegates to ChangeManager."""
        if self._change_manager is None:
            return None
        return self._change_manager.undo()

    def redo(self) -> Any:
        """Redo the most recently undone change. Delegates to ChangeManager."""
        if self._change_manager is None:
            return None
        return self._change_manager.redo()


# Internal helpers used by the methods above
def _shape_to_element(shape: Any, z_order: int) -> Element | None:
    """Convert a single python-pptx *shape* to an :class:`Element`.

    Returns ``None`` for shapes that should be skipped (e.g. empty text shapes).
    """
    etype = _classify_shape(shape)
    position = _build_position(shape)
    element_id = f"{etype}_{uuid.uuid4().hex[:8]}"

    # text logic
    if etype == "text":
        content = _extract_text(shape)
        if not content.strip():
            return None  # skip truly empty text shapes
        fmt = _extract_text_format(shape)
        return Element(
            id=element_id,
            type="text",
            position=position,
            content=content,
            format=fmt,
            z_order=z_order,
        )

    # image logic
    if etype == "image":
        return Element(
            id=element_id,
            type="image",
            position=position,
            content=_extract_image_content(shape),
            z_order=z_order,
        )

    # table logic
    if etype == "table":
        return Element(
            id=element_id,
            type="table",
            position=position,
            content=_extract_table_content(shape),
            z_order=z_order,
        )

    # chart logic
    if etype == "chart":
        return Element(
            id=element_id,
            type="chart",
            position=position,
            content=_extract_chart_content(shape),
            z_order=z_order,
        )

    # diagram / SmartArt logic
    if etype == "diagram":
        return Element(
            id=element_id,
            type="diagram",
            position=position,
            content={
                "name": getattr(shape, "name", "diagram"),
                "description": "SmartArt diagram",
            },
            z_order=z_order,
        )

    # media logic
    if etype == "media":
        return Element(
            id=element_id,
            type="media",
            position=position,
            content={
                "media_type": getattr(shape, "media_type", "unknown")
                if hasattr(shape, "media_type")
                else "unknown",
                "filename": getattr(shape, "name", "media"),
            },
            z_order=z_order,
        )

    # formula (OLE) logic
    if etype == "formula":
        return Element(
            id=element_id,
            type="formula",
            position=position,
            content=_extract_text(shape) or "formula",
            z_order=z_order,
        )

    # generic shape logic
    text = _extract_text(shape)
    return Element(
        id=element_id,
        type="shape",
        position=position,
        content=text if text.strip() else getattr(shape, "name", "shape"),
        z_order=z_order,
    )


def _write_text_element(
    slide: Any,
    element: Element,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    """Add a text box for *element* to *slide* (python-pptx Slide object)."""
    from pptx.util import Pt

    textbox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = textbox.text_frame
    tf.word_wrap = True

    content = element.content if isinstance(element.content, str) else str(element.content)

    # The text frame always has at least one paragraph.
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = content

    if element.format is not None:
        fmt = element.format
        run.font.name = fmt.font_name
        run.font.size = Pt(fmt.font_size)
        run.font.bold = fmt.bold
        run.font.italic = fmt.italic
        run.font.underline = fmt.underline

        # Apply colour
        if fmt.font_color and fmt.font_color != "theme":
            try:
                from pptx.dml.color import RGBColor

                hex_str = fmt.font_color.lstrip("#")
                run.font.color.rgb = RGBColor.from_string(hex_str)
            except Exception:
                pass

        # Alignment (paragraph-level)
        if fmt.alignment:
            from pptx.enum.text import PP_ALIGN

            _align_to_pp = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }
            if fmt.alignment in _align_to_pp:
                p.alignment = _align_to_pp[fmt.alignment]


def _write_image_element(
    slide: Any,
    element: Element,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    """Add a picture for *element* to *slide* (python-pptx Slide object)."""
    if not isinstance(element.content, dict):
        return
    blob = element.content.get("blob")
    if not blob:
        return
    try:
        image_stream = io.BytesIO(blob)
        slide.shapes.add_picture(
            image_stream, int(left), int(top), int(width), int(height)
        )
    except Exception:
        pass


def _write_table_element(
    slide: Any,
    element: Element,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    """Add a table shape for *element* to *slide* (python-pptx Slide object)."""
    rows_data = element.content
    if not isinstance(rows_data, list) or not rows_data:
        return

    rows = len(rows_data)
    cols = max((len(r) for r in rows_data), default=0)
    if cols == 0:
        return

    table_shape = slide.shapes.add_table(
        rows, cols, int(left), int(top), int(width), int(height)
    )
    tbl = table_shape.table
    for r in range(rows):
        row_data = rows_data[r]
        for c in range(min(len(row_data), cols)):
            tbl.cell(r, c).text = str(row_data[c])


def _write_chart_element(
    slide: Any,
    element: Element,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    """Add a chart shape for *element* to *slide* (python-pptx Slide object)."""
    if not isinstance(element.content, dict):
        return

    chart_data = element.content
    chart_type_str = chart_data.get("chart_type", "bar")
    categories = chart_data.get("categories", [])
    series_list = chart_data.get("series", [])
    chart_title = chart_data.get("title")

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    _CHART_TYPE_MAP = {  # noqa: N806
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
    }
    chart_type = _CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_frame = slide.shapes.add_chart(
        chart_type, int(left), int(top), int(width), int(height),
    )
    chart = chart_frame.chart

    cd = CategoryChartData()
    cd.categories = categories
    for s in series_list:
        cd.add_series(s.get("name", ""), s.get("values", []))

    chart.replace_data(cd)

    if chart_title and chart.has_title:
        chart.chart_title.text_frame.paragraphs[0].text = chart_title
